"""Render final Markdown into reviewable office-document formats."""
from __future__ import annotations

import html
import json
import subprocess
from pathlib import Path
from typing import Iterable


ROOT = Path(__file__).resolve().parents[2]


def formats_for_request(request: str, artifact_kind: str) -> list[str]:
    """Choose real formats from explicit user need; never rename Markdown."""
    normalized = request.casefold()
    formats = ["html", "docx", "pdf"]
    if artifact_kind == "financial_model_report":
        formats.append("xlsx")
    if any(term in normalized for term in ("ppt", "presentation", "slide", "발표", "슬라이드")):
        formats.append("pptx")
    if any(term in normalized for term in ("hwp", "hwpx", "한글 문서", "한글파일")):
        formats.append("hwpx")
    return formats


def _plain_lines(markdown: str) -> list[str]:
    result = []
    for line in markdown.splitlines():
        cleaned = line.lstrip("#").strip()
        if cleaned:
            result.append(cleaned)
    return result


def render_bundle(markdown_path: Path, formats: Iterable[str]) -> list[Path]:
    markdown = markdown_path.read_text(encoding="utf-8")
    requested = set(formats)
    rendered: list[Path] = []
    if "html" in requested:
        target = markdown_path.with_suffix(".html")
        body = "\n".join(f"<p>{html.escape(line)}</p>" for line in _plain_lines(markdown))
        target.write_text(
            "<!doctype html><html lang=\"ko\"><meta charset=\"utf-8\">"
            "<title>AI Office Deliverable</title><body>" + body + "</body></html>",
            encoding="utf-8",
        )
        rendered.append(target)
    if "docx" in requested:
        try:
            from docx import Document
        except ImportError as error:
            raise RuntimeError("DOCX rendering requires python-docx") from error
        document = Document()
        for line in markdown.splitlines():
            if line.startswith("# "):
                document.add_heading(line[2:].strip(), level=1)
            elif line.startswith("## "):
                document.add_heading(line[3:].strip(), level=2)
            elif line.strip():
                document.add_paragraph(line.strip())
        target = markdown_path.with_suffix(".docx")
        document.save(target)
        rendered.append(target)
    if "pdf" in requested:
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.cidfonts import UnicodeCIDFont
            from reportlab.pdfgen.canvas import Canvas
        except ImportError as error:
            raise RuntimeError("PDF rendering requires reportlab") from error
        target = markdown_path.with_suffix(".pdf")
        font_name = "HYSMyeongJo-Medium"
        pdfmetrics.registerFont(UnicodeCIDFont(font_name))
        canvas = Canvas(str(target), pagesize=A4)
        width, height = A4
        canvas.setFont(font_name, 10)
        y = height - 50
        for line in _plain_lines(markdown):
            for offset in range(0, max(1, len(line)), 90):
                canvas.drawString(45, y, line[offset : offset + 90])
                y -= 15
                if y < 45:
                    canvas.showPage()
                    canvas.setFont(font_name, 10)
                    y = height - 50
        canvas.save()
        rendered.append(target)
    if "xlsx" in requested:
        try:
            from openpyxl import Workbook
        except ImportError as error:
            raise RuntimeError("XLSX rendering requires openpyxl") from error
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Deliverable"
        sheet.append(["Line", "Content"])
        for index, line in enumerate(_plain_lines(markdown), start=1):
            sheet.append([index, line])
        target = markdown_path.with_suffix(".xlsx")
        workbook.save(target)
        rendered.append(target)
    if "pptx" in requested:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError as error:
            raise RuntimeError("PPTX rendering requires python-pptx") from error
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        title = next(iter(_plain_lines(markdown)), "AI Office Deliverable")
        slide.shapes.title.text = title[:160]
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()
        for index, line in enumerate(_plain_lines(markdown)[1:30]):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = line[:500]
            paragraph.font.size = Pt(14)
        target = markdown_path.with_suffix(".pptx")
        presentation.save(target)
        rendered.append(target)
    if "hwpx" in requested:
        target = markdown_path.with_suffix(".hwpx")
        script = (
            "import { readFileSync, writeFileSync } from 'node:fs';"
            "import { markdownToHwpx, parse } from 'kordoc';"
            "const input = process.argv[1], output = process.argv[2];"
            "const data = await markdownToHwpx(readFileSync(input, 'utf8'));"
            "writeFileSync(output, Buffer.from(data));"
            "const parsed = await parse(output);"
            "if (!parsed?.success || !parsed.markdown?.trim()) process.exit(2);"
        )
        try:
            run = subprocess.run(
                ["node", "--input-type=module", "-e", script, str(markdown_path), str(target)],
                cwd=ROOT,
                capture_output=True,
                text=True,
                timeout=90,
                check=False,
            )
        except (OSError, subprocess.TimeoutExpired) as error:
            raise RuntimeError("HWPX rendering requires local Node.js and kordoc") from error
        if run.returncode:
            raise RuntimeError(f"HWPX rendering/validation failed: {(run.stderr or run.stdout)[:800]}")
        rendered.append(target)
    _validate_rendered(rendered)
    manifest = markdown_path.with_name("ARTIFACTS.json")
    manifest.write_text(
        json.dumps(
            {
                "source": markdown_path.name,
                "rendered": [path.name for path in rendered],
            },
            ensure_ascii=False,
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )
    rendered.append(manifest)
    return rendered


def _validate_rendered(paths: Iterable[Path]) -> None:
    """Re-open generated binary documents; extension alone never counts as evidence."""
    for path in paths:
        if not path.is_file() or path.stat().st_size == 0:
            raise RuntimeError(f"Rendered artifact is missing or empty: {path.name}")
        suffix = path.suffix.lower()
        if suffix == ".docx":
            from docx import Document
            if not "\n".join(paragraph.text for paragraph in Document(path).paragraphs).strip():
                raise RuntimeError("DOCX validation found no readable text")
        elif suffix == ".pdf":
            from pypdf import PdfReader
            if not "".join(page.extract_text() or "" for page in PdfReader(path).pages).strip():
                raise RuntimeError("PDF validation found no extractable text")
        elif suffix == ".xlsx":
            from openpyxl import load_workbook
            workbook = load_workbook(path, read_only=True)
            try:
                rows = workbook.active.max_row
            finally:
                workbook.close()
            if rows < 2:
                raise RuntimeError("XLSX validation found no content rows")
        elif suffix == ".pptx":
            from pptx import Presentation
            if not any(shape.has_text_frame and shape.text.strip() for slide in Presentation(path).slides for shape in slide.shapes):
                raise RuntimeError("PPTX validation found no readable text")
